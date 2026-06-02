"""Generate SRE E-Commerce Dashboard presentation (.pptx)."""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(15, 15, 18)       # charcoal-950
CARD = RGBColor(30, 30, 38)
ACCENT_VIOLET = RGBColor(139, 92, 246)
ACCENT_BLUE = RGBColor(59, 130, 246)
TEXT_LIGHT = RGBColor(226, 232, 240)
TEXT_MUTED = RGBColor(148, 163, 184)
GREEN = RGBColor(34, 197, 94)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(239, 68, 68)
WHITE = RGBColor(255, 255, 255)


def set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str, subtitle: Optional[str] = None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(4)
    # accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(2.2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_VIOLET
    line.line.fill.background()


def add_bullets(slide, items: list[str], left=0.55, top=1.45, width=12.0, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)
        p.bullet = True


def add_footer(slide, text: str = "SRE E-Commerce Dashboard | SRE Cloud"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED


def add_note(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_box_diagram(slide, boxes):
    """boxes: (label, left, top, width, height, fill_color)"""
    for label, l, t, w, h, color in boxes:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or CARD
        shape.line.color.rgb = ACCENT_VIOLET
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER


def add_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = ACCENT_BLUE
    conn.line.width = Pt(2)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # gradient bar top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_VIOLET
    bar.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "SRE E-Commerce Dashboard"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(11.5), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "Real-time observability for a microservices e-commerce stack"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_MUTED

    tag = slide.shapes.add_textbox(Inches(0.8), Inches(4.35), Inches(11.0), Inches(0.6))
    p3 = tag.text_frame.paragraphs[0]
    p3.text = "Simulate, observe, break, and heal — like production SRE"
    p3.font.size = Pt(18)
    p3.font.color.rgb = ACCENT_BLUE

    chips = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5))
    pc = chips.text_frame.paragraphs[0]
    pc.text = "React  •  FastAPI  •  Docker  •  Prometheus  •  Grafana  •  Jaeger"
    pc.font.size = Pt(14)
    pc.font.color.rgb = GREEN

    # decorative nodes
    for x, y, c in [(9.5, 2.5, ACCENT_VIOLET), (10.5, 3.8, ACCENT_BLUE), (9.0, 5.0, GREEN)]:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
        circ.fill.solid()
        circ.fill.fore_color.rgb = c
        circ.line.fill.background()

    add_footer(slide, "SRE Cloud Project | Portfolio Demo")
    add_note(
        slide,
        "Introduce as production-style microservices stack with SRE dashboard. "
        "One command docker compose up brings full stack.",
    )


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Why This Project?")
    add_bullets(
        slide,
        [
            "Microservices fail in subtle ways — one slow DB affects the entire user journey",
            "SRE teams need one pane of glass: health, metrics, incidents, and experiments",
            "Learn chaos engineering and load testing safely before production",
            "Bridges the gap between toy apps and real production observability",
            "Demonstrates live data provenance — not just mocked charts",
        ],
        top=1.5,
    )
    # before/after boxes
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.0), Inches(2.5), Inches(2.2))
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(60, 20, 20)
    b1.text_frame.paragraphs[0].text = "Before\nFragmented\nalerts"
    b1.text_frame.paragraphs[0].font.color.rgb = RED
    b1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(2.0), Inches(2.5), Inches(2.2))
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(20, 50, 30)
    b2.text_frame.paragraphs[0].text = "After\nUnified SRE\nDashboard"
    b2.text_frame.paragraphs[0].font.color.rgb = GREEN
    b2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_footer(slide)
    add_note(slide, "Frame pain: alert fatigue, unclear if metrics are real, hard to demo SRE in interviews.")


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "What We Built")
    add_bullets(
        slide,
        [
            "3-tier microservices e-commerce: product, cart, order services",
            "API Gateway with live HTTP probes + chaos overlays + incident simulation",
            "React dashboard: Overview, Microservices, Incidents, Chaos & Load",
            "Full observability via Docker: Prometheus, Grafana, Alertmanager, Jaeger",
            "PostgreSQL 16 + Redis 7 for realistic data-layer dependencies",
        ],
        top=1.5,
        width=6.5,
    )
    add_box_diagram(
        slide,
        [
            ("User", 7.5, 2.0, 1.4, 0.55, CARD),
            ("Dashboard", 7.2, 3.0, 1.8, 0.55, ACCENT_VIOLET),
            ("API Gateway", 7.0, 4.0, 2.2, 0.55, ACCENT_BLUE),
            ("Services x3", 7.0, 5.0, 2.2, 0.55, CARD),
            ("DB + Redis", 7.0, 6.0, 2.2, 0.55, CARD),
        ],
    )
    add_footer(slide)
    add_note(slide, "docker compose up --build -d starts 10+ containers on ecommerce-net bridge.")


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "High-Level Architecture", "Docker Compose — ecommerce-net")
    boxes = [
        ("Browser :5173", 0.6, 1.6, 2.0, 0.5, CARD),
        ("React Dashboard", 0.5, 2.4, 2.2, 0.5, ACCENT_VIOLET),
        ("API Gateway :8080", 0.4, 3.3, 2.4, 0.55, ACCENT_BLUE),
        ("product :5001", 3.2, 4.2, 1.8, 0.5, CARD),
        ("cart :5002", 5.3, 4.2, 1.6, 0.5, CARD),
        ("order :5003", 7.2, 4.2, 1.6, 0.5, CARD),
        ("PostgreSQL", 3.0, 5.3, 2.0, 0.5, CARD),
        ("Redis", 5.5, 5.3, 1.5, 0.5, CARD),
        ("Prometheus :9090", 9.0, 2.0, 2.2, 0.5, CARD),
        ("Grafana :3000", 9.0, 2.8, 2.2, 0.5, CARD),
        ("Alertmanager", 9.0, 3.6, 2.2, 0.5, CARD),
        ("Jaeger :16686", 9.0, 4.4, 2.2, 0.5, CARD),
    ]
    add_box_diagram(slide, boxes)
    add_bullets(
        slide,
        [
            "10+ containers with healthchecks on Postgres & Redis",
            "Cart service calls product-service; all probed by gateway",
            "Prometheus scrapes /metrics on all services every 15s",
        ],
        left=0.55,
        top=6.0,
        width=12.0,
        height=1.2,
        size=14,
    )
    add_footer(slide)
    add_note(
        slide,
        "api-gateway depends on all microservices. Vite dev server proxies /api to :8080.",
    )


def slide_microservices(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "E-Commerce Microservices")
    table = slide.shapes.add_table(4, 4, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.0)).table
    headers = ["Service", "Port", "Responsibility", "Dependencies"]
    rows = [
        ["product-service", "5001", "Catalog SKUs, /health, /metrics", "PostgreSQL, Redis"],
        ["cart-service", "5002", "Cart operations", "Redis, product-service"],
        ["order-service", "5003", "Order placement", "PostgreSQL, Redis"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_VIOLET
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_LIGHT
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD
    add_bullets(
        slide,
        [
            "FastAPI + shared ServiceMetrics middleware on each service",
            "Endpoints: /health, /, /metrics (Prometheus exposition)",
            "Built from services/*/Dockerfile with services/shared/metrics.py",
        ],
        top=3.7,
        size=16,
    )
    add_footer(slide)
    add_note(slide, "Env vars: DATABASE_URL, REDIS_URL, PRODUCT_SERVICE_URL per service.")


def slide_gateway(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "API Gateway — Brain of the Stack", "backend/main.py + service_probe.py")
    add_bullets(
        slide,
        [
            "Async HTTP probes merge live status with chaos simulation overlays",
            "Chaos: kill service, slow DB, high latency — auto-creates incidents",
            "Aggregates live metrics when services up; simulates when offline",
            "Interactive OpenAPI docs at http://localhost:8080/docs",
        ],
        top=1.45,
        width=5.8,
        size=16,
    )
    api_box = slide.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(6.2), Inches(5.0))
    tf = api_box.text_frame
    tf.word_wrap = True
    lines = [
        "GET  /health, /metrics",
        "GET  /services/status, /services/detail",
        "GET  /incidents",
        "POST /incidents/{id}/resolve",
        "POST /simulate/kill-service",
        "POST /simulate/slow-db",
        "POST /simulate/high-latency",
        "POST /simulate/reset",
        "POST /load-test",
        "GET  /load-test/results",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.font.color.rgb = GREEN
        p.space_after = Pt(4)
    add_footer(slide)
    add_note(slide, "Frontend polls every 3s. httpx async client for service probes.")


def slide_overview_ui(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Dashboard — Overview", "http://localhost:5173")
    add_bullets(
        slide,
        [
            "KPI bento grid: success rate, P95 latency, error rate, active users",
            "SVG LineChart — 24-point latency & error rate history",
            "TrafficMap visualizes inter-service traffic flow",
            "ResourceBar: CPU, memory, storage utilization",
            "Health badge: healthy / degraded / critical — polls every 3s",
        ],
        top=1.45,
        width=5.5,
        size=15,
    )
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.0)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = CARD
    placeholder.line.color.rgb = ACCENT_BLUE
    tf = placeholder.text_frame
    tf.paragraphs[0].text = "[ Insert screenshot:\nOverview page ]\n\nRun: npm run dev\nCapture localhost:5173"
    for p in tf.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER
    add_footer(slide)
    add_note(slide, "React 18, TypeScript, Vite, Tailwind, Framer Motion, Lucide icons.")


def slide_services_ui(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Microservices View — Live Provenance", "/services route")
    add_bullets(
        slide,
        [
            "ServiceMesh component shows inter-service topology",
            "Per-service: HTTP status, latency, probe URL, JSON payload",
            "Badges: LIVE | LIVE + CHAOS | UNREACHABLE",
            "GET /services/detail exposes full provenance for demos",
        ],
        top=1.45,
        width=5.5,
    )
    ph = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(1.5), Inches(6.5), Inches(5.0)
    )
    ph.fill.solid()
    ph.fill.fore_color.rgb = CARD
    ph.line.color.rgb = GREEN
    ph.text_frame.paragraphs[0].text = "[ Insert screenshot:\nMicroservices page\nwith LIVE badges ]"
    ph.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED
    ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_footer(slide)
    add_note(slide, "Proves real HTTP probes to ports 5001-5003 — not fake data.")


def slide_incidents(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Incident Management")
    add_bullets(
        slide,
        [
            "Active vs Resolved columns with severity: Critical, Warning, Info",
            "Incidents auto-created on chaos events (kill, slow DB, latency)",
            "One-click Resolve → POST /incidents/{id}/resolve",
            "AnimatePresence UI — mirrors PagerDuty-style on-call workflow",
        ],
        top=1.5,
    )
    for i, (sev, color, x) in enumerate(
        [("Critical", RED, 7.0), ("Warning", AMBER, 9.2), ("Info", ACCENT_BLUE, 11.4)]
    ):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.5), Inches(2.0), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = color
        tf = card.text_frame
        tf.paragraphs[0].text = sev
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].font.bold = True
        tf.add_paragraph().text = "Sample incident"
        tf.paragraphs[1].font.size = Pt(10)
        tf.paragraphs[1].font.color.rgb = TEXT_MUTED
    add_footer(slide)
    add_note(slide, "Run chaos first to populate active incidents for live demo.")


def slide_chaos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Chaos Engineering & Load Testing", "/chaos route")
    add_bullets(
        slide,
        [
            "Chaos: Kill Service, Slow DB, High Latency, Reset All",
            "Load test: duration, concurrency, target (gateway or any service)",
            "Terminal-style log output in the UI",
            "Real-time chart: RPS, error %, avg latency ms",
            "Workflow: break → observe incidents → load test → reset/heal",
        ],
        top=1.45,
        width=6.0,
    )
    chaos = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.5), Inches(3.0), Inches(4.8)
    )
    chaos.fill.solid()
    chaos.fill.fore_color.rgb = RGBColor(50, 20, 20)
    chaos.text_frame.paragraphs[0].text = "CHAOS\nKill Service\nSlow DB\nHigh Latency\nReset"
    chaos.text_frame.paragraphs[0].font.color.rgb = RED

    load = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.7), Inches(1.5), Inches(3.0), Inches(4.8)
    )
    load.fill.solid()
    load.fill.fore_color.rgb = RGBColor(20, 40, 50)
    load.text_frame.paragraphs[0].text = "LOAD TEST\nTerminal logs\nRPS chart\nPOST /load-test"
    load.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    add_footer(slide)
    add_note(slide, "POST /simulate/* and POST /load-test; results at GET /load-test/results.")


def slide_observability(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Observability Stack")
    tools = [
        ("Prometheus", ":9090", "Metrics TSDB, scrapes all services"),
        ("Grafana", ":3000", "Dashboards (admin/admin)"),
        ("Alertmanager", ":9093", "Alert routing"),
        ("Jaeger", ":16686", "Distributed tracing UI"),
    ]
    for i, (name, port, desc) in enumerate(tools):
        x = 0.5 + (i % 2) * 6.3
        y = 1.5 + (i // 2) * 2.5
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.0), Inches(2.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = ACCENT_VIOLET
        tf = card.text_frame
        tf.paragraphs[0].text = f"{name} {port}"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_MUTED
    add_bullets(
        slide,
        ["prometheus.yml jobs: api-gateway, product, cart, order (15s scrape)"],
        top=6.2,
        size=14,
    )
    add_footer(slide)
    add_note(slide, "Grafana datasource pre-provisioned to Prometheus in observability/grafana/.")


def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Tech Stack & Project Structure")
    add_bullets(
        slide,
        [
            "Frontend: React 18, TypeScript, Vite, Tailwind, Framer Motion",
            "Backend: FastAPI, Uvicorn, httpx, Pydantic",
            "Infra: Docker Compose, PostgreSQL 16, Redis 7",
            "Observability: Prometheus, Grafana, Alertmanager, Jaeger",
        ],
        top=1.45,
        width=5.8,
        size=15,
    )
    tree = slide.shapes.add_textbox(Inches(6.4), Inches(1.45), Inches(6.4), Inches(5.2))
    tf = tree.text_frame
    structure = """SRE Cloud/
├── docker-compose.yml
├── backend/           API gateway
├── services/          product, cart, order
│   └── shared/        metrics.py
├── frontend/src/
│   ├── views/         Overview, Services, Incidents, Chaos
│   └── components/    Charts, TrafficMap, ServiceMesh
├── observability/
│   ├── prometheus/
│   ├── grafana/
│   └── alertmanager/
└── scripts/           start-microservices.ps1"""
    for i, line in enumerate(structure.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = GREEN
    add_footer(slide)
    add_note(slide, "Local dev: backend :8080, frontend :5173, Vite proxy /api → 8080.")


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Demo Flow & Future Work")
    add_bullets(
        slide,
        [
            "1. docker compose up --build -d",
            "2. Open http://localhost:5173 (npm run dev for frontend)",
            "3. Demo: Overview healthy → Microservices LIVE → Chaos kill → Incidents → Load test → Reset",
            "Future: Kubernetes/Helm, OpenTelemetry → Jaeger, CI/CD synthetic monitoring",
        ],
        top=1.5,
        size=17,
    )
    outcomes = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.8)
    )
    outcomes.fill.solid()
    outcomes.fill.fore_color.rgb = CARD
    outcomes.line.color.rgb = GREEN
    tf = outcomes.text_frame
    tf.paragraphs[0].text = "Key outcomes: End-to-end microservices + SRE workflows | Live probes with simulation fallback | Interview-ready architecture narrative"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = TEXT_LIGHT

    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6))
    thanks.text_frame.paragraphs[0].text = "Thank you — Questions?"
    thanks.text_frame.paragraphs[0].font.size = Pt(28)
    thanks.text_frame.paragraphs[0].font.bold = True
    thanks.text_frame.paragraphs[0].font.color.rgb = ACCENT_VIOLET
    thanks.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_footer(slide, "github.com — SRE Cloud | Add your name & contact")
    add_note(
        slide,
        "Swagger at :8080/docs. Grafana admin/admin. Optional live 2-min demo.",
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_microservices(prs)
    slide_gateway(prs)
    slide_overview_ui(prs)
    slide_services_ui(prs)
    slide_incidents(prs)
    slide_chaos(prs)
    slide_observability(prs)
    slide_tech(prs)
    slide_demo(prs)

    out = Path(__file__).resolve().parent.parent / "SRE_E-Commerce_Dashboard.pptx"
    prs.save(str(out))
    print(f"Created: {out}")


if __name__ == "__main__":
    main()
